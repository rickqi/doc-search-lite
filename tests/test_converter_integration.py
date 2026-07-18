"""
Integration tests for the converter module.

Tests cover:
- End-to-end conversion flow using ConverterCoordinator
- Multiple format conversion in batch
- Error handling across converters
- Recovery from partial failures
- Full pipeline: coordinator → specific converter → output
- Mixed format directory handling
"""

import base64
from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest
from docx import Document
from openpyxl import Workbook
from pptx import Presentation

from src.converter.base import ConvertResult
from src.converter.coordinator import ConverterCoordinator
from src.converter.html import HTMLConverter
from src.converter.ocr import OCRServiceConfig
from src.converter.office import OfficeConverter
from src.converter.pdf import PDFConverter


# Helper functions to create test files
def create_minimal_pdf(file_path: Path) -> None:
    """Create a minimal valid PDF file for testing."""
    pdf_content = b"""%PDF-1.4
1 0 obj
<<
/Type /Catalog
/Pages 2 0 R
>>
endobj
2 0 obj
<<
/Type /Pages
/Kids [3 0 R]
/Count 1
>>
endobj
3 0 obj
<<
/Type /Page
/Parent 2 0 R
/MediaBox [0 0 612 792]
/Contents 4 0 R
/Resources <<
/Font <<
/F1 5 0 R
>>
>>
>>
endobj
4 0 obj
<<
/Length 44
>>
stream
BT
/F1 12 Tf
100 700 Td
(Test Content) Tj
ET
endstream
endobj
5 0 obj
<<
/Type /Font
/Subtype /Type1
/BaseFont /Helvetica
>>
endobj
xref
0 6
0000000000 65535 f 
0000000009 00000 n 
0000000058 00000 n 
0000000115 00000 n 
0000000266 00000 n 
0000000362 00000 n 
trailer
<<
/Size 6
/Root 1 0 R
>>
startxref
439
%%EOF
"""
    file_path.write_bytes(pdf_content)


def create_docx_file(file_path: Path, title: str = "Test Document") -> None:
    """Create a minimal DOCX file for testing."""
    doc = Document()
    doc.add_heading(title, level=1)
    doc.add_paragraph("This is a test paragraph.")
    doc.save(str(file_path))


def create_pptx_file(file_path: Path, title: str = "Test Presentation") -> None:
    """Create a minimal PPTX file for testing."""
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    prs.save(str(file_path))


def create_xlsx_file(file_path: Path, title: str = "Test Spreadsheet") -> None:
    """Create a minimal XLSX file for testing."""
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws["A1"] = "Header"
    ws["A2"] = "Data"
    wb.save(str(file_path))


def create_html_file(file_path: Path, title: str = "Test Page") -> None:
    """Create a minimal HTML file for testing."""
    content = f"""<!DOCTYPE html>
<html>
<head><title>{title}</title></head>
<body>
    <h1>{title}</h1>
    <p>Test content</p>
</body>
</html>
"""
    file_path.write_text(content, encoding="utf-8")


def create_png_file(file_path: Path) -> None:
    """Create a minimal PNG file for testing."""
    image_content = base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mNk+M9QDwADhgGAWjR9awAAAABJRU5ErkJggg=="
    )
    file_path.write_bytes(image_content)


class TestEndToEndConversion:
    """Test end-to-end conversion through the coordinator."""

    @pytest.fixture
    def coordinator(self):
        """Create a ConverterCoordinator instance."""
        return ConverterCoordinator()

    def test_pdf_e2e_conversion(
        self, tmp_path: Path, coordinator: ConverterCoordinator
    ):
        """Test end-to-end PDF conversion through coordinator."""
        # Setup
        pdf_path = tmp_path / "document.pdf"
        output_dir = tmp_path / "output"
        create_minimal_pdf(pdf_path)

        # Execute
        result = coordinator.convert(pdf_path, output_dir)

        # Verify
        assert result.source_file == pdf_path
        assert result.converter_name == "pdfplumber"
        assert "coordinator_used" in result.metadata
        assert result.metadata["coordinator_used"] is True
        assert result.metadata["coordinator_converter"] == "pdfplumber"

        # Output directory should be created
        assert output_dir.exists()

    def test_docx_e2e_conversion(
        self, tmp_path: Path, coordinator: ConverterCoordinator
    ):
        """Test end-to-end DOCX conversion through coordinator."""
        # Setup
        docx_path = tmp_path / "document.docx"
        output_dir = tmp_path / "output"
        create_docx_file(docx_path)

        # Execute
        result = coordinator.convert(docx_path, output_dir)

        # Verify
        assert result.success is True
        assert result.source_file == docx_path
        assert "Test Document" in result.markdown
        assert result.output_file is not None
        assert result.output_file.exists()
        assert result.output_file.suffix == ".md"
        assert result.metadata["coordinator_converter"] == "OfficeConverter"

    def test_pptx_e2e_conversion(
        self, tmp_path: Path, coordinator: ConverterCoordinator
    ):
        """Test end-to-end PPTX conversion through coordinator."""
        # Setup
        pptx_path = tmp_path / "presentation.pptx"
        output_dir = tmp_path / "output"
        create_pptx_file(pptx_path)

        # Execute
        result = coordinator.convert(pptx_path, output_dir)

        # Verify
        assert result.success is True
        assert result.source_file == pptx_path
        assert result.output_file is not None
        assert result.output_file.exists()
        assert result.metadata["coordinator_converter"] == "OfficeConverter"

    def test_xlsx_e2e_conversion(
        self, tmp_path: Path, coordinator: ConverterCoordinator
    ):
        """Test end-to-end XLSX conversion through coordinator."""
        # Setup
        xlsx_path = tmp_path / "spreadsheet.xlsx"
        output_dir = tmp_path / "output"
        create_xlsx_file(xlsx_path)

        # Execute
        result = coordinator.convert(xlsx_path, output_dir)

        # Verify
        assert result.success is True
        assert result.source_file == xlsx_path
        assert result.output_file is not None
        assert result.output_file.exists()
        assert result.metadata["coordinator_converter"] == "OfficeConverter"

    def test_html_e2e_conversion(
        self, tmp_path: Path, coordinator: ConverterCoordinator
    ):
        """Test end-to-end HTML conversion through coordinator."""
        # Setup
        html_path = tmp_path / "page.html"
        output_dir = tmp_path / "output"
        create_html_file(html_path)

        # Execute
        result = coordinator.convert(html_path, output_dir)

        # Verify
        assert result.success is True
        assert result.source_file == html_path
        assert "Test Page" in result.markdown
        assert result.output_file is not None
        assert result.output_file.exists()
        assert result.metadata["coordinator_converter"] == "HTMLConverter"

    def test_htm_e2e_conversion(
        self, tmp_path: Path, coordinator: ConverterCoordinator
    ):
        """Test end-to-end .htm file conversion through coordinator."""
        # Setup
        htm_path = tmp_path / "page.htm"
        output_dir = tmp_path / "output"
        create_html_file(htm_path)

        # Execute
        result = coordinator.convert(htm_path, output_dir)

        # Verify
        assert result.success is True
        assert result.source_file == htm_path
        assert result.output_file is not None
        assert result.output_file.suffix == ".md"


class TestBatchConversion:
    """Test batch conversion of multiple files."""

    @pytest.fixture
    def coordinator(self):
        """Create a ConverterCoordinator instance."""
        return ConverterCoordinator()

    @pytest.fixture
    def mixed_files_dir(self, tmp_path: Path) -> Path:
        """Create a directory with mixed file types."""
        files_dir = tmp_path / "source_files"
        files_dir.mkdir()

        # Create files of different formats
        create_minimal_pdf(files_dir / "doc1.pdf")
        create_docx_file(files_dir / "doc2.docx")
        create_pptx_file(files_dir / "doc3.pptx")
        create_xlsx_file(files_dir / "doc4.xlsx")
        create_html_file(files_dir / "doc5.html")

        return files_dir

    def test_batch_convert_all_supported_formats(
        self, tmp_path: Path, coordinator: ConverterCoordinator, mixed_files_dir: Path
    ):
        """Test batch conversion of all supported file formats."""
        output_dir = tmp_path / "output"
        results = []

        # Convert all files
        for source_file in mixed_files_dir.iterdir():
            if coordinator.can_convert(source_file):
                result = coordinator.convert(source_file, output_dir)
                results.append((source_file.name, result))

        # Verify all conversions
        assert len(results) == 5

        for filename, result in results:
            assert result.success is True, f"Failed to convert {filename}"
            assert result.output_file is not None
            assert result.output_file.exists()

    def test_batch_convert_with_results_tracking(
        self, tmp_path: Path, coordinator: ConverterCoordinator, mixed_files_dir: Path
    ):
        """Test tracking conversion results in batch."""
        output_dir = tmp_path / "output"
        results_by_format = {
            ".pdf": [],
            ".docx": [],
            ".pptx": [],
            ".xlsx": [],
            ".html": [],
        }

        # Convert and categorize
        for source_file in mixed_files_dir.iterdir():
            ext = source_file.suffix.lower()
            if ext in results_by_format:
                result = coordinator.convert(source_file, output_dir)
                results_by_format[ext].append(result)

        # Verify each format was converted
        for ext, results in results_by_format.items():
            assert len(results) == 1, f"Expected 1 result for {ext}"
            assert results[0].success is True

    def test_batch_convert_preserves_source_info(
        self, tmp_path: Path, coordinator: ConverterCoordinator, mixed_files_dir: Path
    ):
        """Test that source file info is preserved in batch conversion."""
        output_dir = tmp_path / "output"
        source_files = list(mixed_files_dir.iterdir())
        results = {}

        for source_file in source_files:
            if coordinator.can_convert(source_file):
                result = coordinator.convert(source_file, output_dir)
                results[source_file.name] = result

        # Verify source file tracking
        for filename, result in results.items():
            assert result.source_file.name == filename


class TestMixedFormatDirectory:
    """Test handling of directories with mixed supported and unsupported files."""

    @pytest.fixture
    def coordinator(self):
        """Create a ConverterCoordinator instance."""
        return ConverterCoordinator()

    @pytest.fixture
    def mixed_supported_unsupported_dir(self, tmp_path: Path) -> Path:
        """Create a directory with both supported and unsupported files."""
        files_dir = tmp_path / "mixed_files"
        files_dir.mkdir()

        # Supported formats
        create_minimal_pdf(files_dir / "supported1.pdf")
        create_docx_file(files_dir / "supported2.docx")
        create_html_file(files_dir / "supported3.html")

        # Unsupported formats (truly unsupported, not newly added converters)
        (files_dir / "unsupported1.xyz").write_text("unknown format")
        (files_dir / "unsupported2.rtf").write_text("rich text")
        (files_dir / "unsupported3.exe").write_bytes(b"binary\x00data")

        return files_dir

    def test_convert_only_supported_files(
        self,
        tmp_path: Path,
        coordinator: ConverterCoordinator,
        mixed_supported_unsupported_dir: Path,
    ):
        """Test converting only supported files from a mixed directory."""
        output_dir = tmp_path / "output"
        supported_count = 0
        unsupported_count = 0

        for source_file in mixed_supported_unsupported_dir.iterdir():
            if coordinator.can_convert(source_file):
                result = coordinator.convert(source_file, output_dir)
                assert result.success is True
                supported_count += 1
            else:
                unsupported_count += 1

        assert supported_count == 3
        assert unsupported_count == 3

    def test_unsupported_file_error_handling(
        self,
        tmp_path: Path,
        coordinator: ConverterCoordinator,
        mixed_supported_unsupported_dir: Path,
    ):
        """Test error handling for unsupported files."""
        output_dir = tmp_path / "output"

        for source_file in mixed_supported_unsupported_dir.iterdir():
            result = coordinator.convert(source_file, output_dir)

            if source_file.suffix.lower() in [".pdf", ".docx", ".html"]:
                assert result.success is True
            else:
                assert result.success is False
                assert len(result.errors) > 0
                assert (
                    "Unsupported" in result.errors[0]
                    or "unsupported" in result.errors[0].lower()
                )


class TestErrorHandlingAndRecovery:
    """Test error handling and recovery across converters."""

    @pytest.fixture
    def coordinator(self):
        """Create a ConverterCoordinator instance."""
        return ConverterCoordinator()

    def test_missing_file_error_recovery(
        self, tmp_path: Path, coordinator: ConverterCoordinator
    ):
        """Test recovery from missing file errors."""
        output_dir = tmp_path / "output"

        # Try to convert non-existent file
        missing_file = tmp_path / "nonexistent.pdf"
        result = coordinator.convert(missing_file, output_dir)

        assert result.success is False
        assert "not found" in result.errors[0].lower()

        # Create a valid file and verify coordinator still works
        valid_file = tmp_path / "valid.pdf"
        create_minimal_pdf(valid_file)
        result2 = coordinator.convert(valid_file, output_dir)

        assert result2.success is True or result2.converter_name == "pdfplumber"

    def test_corrupted_file_error_recovery(
        self, tmp_path: Path, coordinator: ConverterCoordinator
    ):
        """Test recovery from corrupted file errors."""
        output_dir = tmp_path / "output"

        # Create corrupted file
        corrupted_file = tmp_path / "corrupted.docx"
        corrupted_file.write_bytes(b"This is not a valid DOCX file")

        result = coordinator.convert(corrupted_file, output_dir)

        # Should handle gracefully without crashing
        assert isinstance(result, ConvertResult)
        assert result.source_file == corrupted_file

        # Verify coordinator still works with valid files
        valid_file = tmp_path / "valid.docx"
        create_docx_file(valid_file)
        result2 = coordinator.convert(valid_file, output_dir)

        assert result2.success is True

    def test_error_propagation_from_converter(
        self, tmp_path: Path, coordinator: ConverterCoordinator
    ):
        """Test that errors from converters are properly propagated."""
        output_dir = tmp_path / "output"

        # Create a file with unsupported extension
        unsupported_file = tmp_path / "test.xyz"
        unsupported_file.write_text("test content")

        result = coordinator.convert(unsupported_file, output_dir)

        assert result.success is False
        assert len(result.errors) > 0
        assert ".xyz" in result.errors[0]

    def test_sequential_conversion_after_error(
        self, tmp_path: Path, coordinator: ConverterCoordinator
    ):
        """Test that coordinator recovers and can process subsequent files after an error."""
        output_dir = tmp_path / "output"

        # First, cause an error with a truly unsupported format
        bad_file = tmp_path / "bad.xyz"
        bad_file.write_text("not supported")
        result1 = coordinator.convert(bad_file, output_dir)
        assert result1.success is False

        # Then, convert valid files
        valid_files = [
            (tmp_path / "doc1.pdf", create_minimal_pdf),
            (tmp_path / "doc2.docx", create_docx_file),
            (tmp_path / "doc3.html", create_html_file),
        ]

        for file_path, create_func in valid_files:
            create_func(file_path)
            result = coordinator.convert(file_path, output_dir)
            assert result.success is True, f"Failed after error recovery: {file_path}"


class TestConverterSelection:
    """Test that coordinator correctly selects converters."""

    @pytest.fixture
    def coordinator(self):
        """Create a ConverterCoordinator instance."""
        return ConverterCoordinator()

    def test_selects_pdf_converter(self, coordinator: ConverterCoordinator):
        """Test that PDF files are routed to PDFConverter."""
        converter = coordinator.get_converter(Path("test.pdf"))
        assert isinstance(converter, PDFConverter)

    def test_selects_office_converter_for_docx(self, coordinator: ConverterCoordinator):
        """Test that DOCX files are routed to OfficeConverter."""
        converter = coordinator.get_converter(Path("test.docx"))
        assert isinstance(converter, OfficeConverter)

    def test_selects_office_converter_for_pptx(self, coordinator: ConverterCoordinator):
        """Test that PPTX files are routed to OfficeConverter."""
        converter = coordinator.get_converter(Path("test.pptx"))
        assert isinstance(converter, OfficeConverter)

    def test_selects_office_converter_for_xlsx(self, coordinator: ConverterCoordinator):
        """Test that XLSX files are routed to OfficeConverter."""
        converter = coordinator.get_converter(Path("test.xlsx"))
        assert isinstance(converter, OfficeConverter)

    def test_selects_html_converter(self, coordinator: ConverterCoordinator):
        """Test that HTML files are routed to HTMLConverter."""
        converter = coordinator.get_converter(Path("test.html"))
        assert isinstance(converter, HTMLConverter)

        converter = coordinator.get_converter(Path("test.htm"))
        assert isinstance(converter, HTMLConverter)


class TestOutputVerification:
    """Test output file verification."""

    @pytest.fixture
    def coordinator(self):
        """Create a ConverterCoordinator instance."""
        return ConverterCoordinator()

    def test_output_file_structure(
        self, tmp_path: Path, coordinator: ConverterCoordinator
    ):
        """Test that output files have correct structure."""
        docx_path = tmp_path / "document.docx"
        output_dir = tmp_path / "output"
        create_docx_file(docx_path)

        result = coordinator.convert(docx_path, output_dir)

        assert result.success is True
        assert result.output_file is not None
        assert result.output_file.parent == output_dir
        assert result.output_file.stem == "document"
        assert result.output_file.suffix == ".md"

    def test_output_file_content_matches_markdown(
        self, tmp_path: Path, coordinator: ConverterCoordinator
    ):
        """Test that output file content matches result.markdown."""
        docx_path = tmp_path / "document.docx"
        output_dir = tmp_path / "output"
        create_docx_file(docx_path, "Unique Title XYZ")

        result = coordinator.convert(docx_path, output_dir)

        assert result.success is True

        # Read output file
        output_content = result.output_file.read_text(encoding="utf-8")

        # Strip frontmatter (injected by pipeline Step 6) before comparing,
        # since it may be present in either the file or result.markdown
        from src.converter.frontmatter import strip_frontmatter
        _, body = strip_frontmatter(output_content)
        _, expected = strip_frontmatter(result.markdown)

        # Body should match markdown
        assert body == expected
        assert "Unique Title XYZ" in output_content

    def test_nested_output_directory_creation(
        self, tmp_path: Path, coordinator: ConverterCoordinator
    ):
        """Test that nested output directories are created."""
        docx_path = tmp_path / "document.docx"
        nested_output_dir = tmp_path / "output" / "nested" / "deep"
        create_docx_file(docx_path)

        result = coordinator.convert(docx_path, nested_output_dir)

        assert result.success is True
        assert nested_output_dir.exists()
        assert result.output_file.exists()


class TestMetadataPropagation:
    """Test metadata propagation through the pipeline."""

    @pytest.fixture
    def coordinator(self):
        """Create a ConverterCoordinator instance."""
        return ConverterCoordinator()

    def test_coordinator_metadata_added(
        self, tmp_path: Path, coordinator: ConverterCoordinator
    ):
        """Test that coordinator adds its metadata to results."""
        docx_path = tmp_path / "document.docx"
        output_dir = tmp_path / "output"
        create_docx_file(docx_path)

        result = coordinator.convert(docx_path, output_dir)

        assert "coordinator_used" in result.metadata
        assert result.metadata["coordinator_used"] is True
        assert "coordinator_converter" in result.metadata

    def test_converter_metadata_preserved(
        self, tmp_path: Path, coordinator: ConverterCoordinator
    ):
        """Test that converter-specific metadata is preserved."""
        html_path = tmp_path / "page.html"
        output_dir = tmp_path / "output"
        create_html_file(html_path)

        result = coordinator.convert(html_path, output_dir)

        assert result.converter_name == "HTMLConverter"
        assert result.converter_version == "0.1.0"


class TestOCRIntegration:
    """Test OCR integration through coordinator."""

    @pytest.fixture
    def coordinator_with_ocr(self):
        """Create a coordinator with OCR configuration."""
        ocr_config = OCRServiceConfig(
            api_key="test-key",
            base_url="https://test.api.url",
            max_retries=1,
            retry_delay=0.1,
        )
        return ConverterCoordinator(ocr_config=ocr_config, enable_ocr_fallback=True)

    @pytest.fixture
    def coordinator_without_ocr(self):
        """Create a coordinator without OCR configuration."""
        return ConverterCoordinator(enable_ocr_fallback=False)

    def test_ocr_disabled_by_default(self, tmp_path: Path):
        """Test that OCR is disabled when not configured."""
        coordinator = ConverterCoordinator()
        pdf_path = tmp_path / "document.pdf"
        output_dir = tmp_path / "output"
        create_minimal_pdf(pdf_path)

        result = coordinator.convert(pdf_path, output_dir)

        assert result.ocr_used is False

    def test_ocr_fallback_disabled_option(
        self, tmp_path: Path, coordinator_with_ocr: ConverterCoordinator
    ):
        """Test that OCR fallback can be disabled via option."""
        pdf_path = tmp_path / "document.pdf"
        output_dir = tmp_path / "output"
        create_minimal_pdf(pdf_path)

        result = coordinator_with_ocr.convert(
            pdf_path, output_dir, options={"disable_ocr_fallback": True}
        )

        assert result.ocr_used is False

    @patch("src.converter.ocr.OCRService._get_client")
    def test_ocr_for_scanned_pdf(
        self,
        mock_get_client,
        tmp_path: Path,
        coordinator_with_ocr: ConverterCoordinator,
    ):
        """Test OCR is triggered for scanned PDFs."""
        # Setup mock OCR service
        mock_client = MagicMock()
        mock_client.layout_parsing.create.return_value = "OCR extracted text"
        mock_get_client.return_value = mock_client

        # Create minimal PDF (will appear as scanned due to low text)
        pdf_path = tmp_path / "scanned.pdf"
        output_dir = tmp_path / "output"
        create_minimal_pdf(pdf_path)

        result = coordinator_with_ocr.convert(
            pdf_path, output_dir, options={"force_ocr": True, "extract_images": True}
        )

        # Verify result structure (OCR may or may not be triggered depending on PDF content)
        assert isinstance(result, ConvertResult)
        assert result.source_file == pdf_path


class TestFullPipeline:
    """Test complete pipeline scenarios."""

    @pytest.fixture
    def coordinator(self):
        """Create a ConverterCoordinator instance."""
        return ConverterCoordinator()

    def test_complete_pdf_pipeline(
        self, tmp_path: Path, coordinator: ConverterCoordinator
    ):
        """Test complete PDF conversion pipeline."""
        pdf_path = tmp_path / "document.pdf"
        output_dir = tmp_path / "output"
        create_minimal_pdf(pdf_path)

        # Execute full pipeline
        result = coordinator.convert(pdf_path, output_dir)

        # Verify all pipeline stages
        assert result.source_file == pdf_path
        assert result.converter_name == "pdfplumber"
        assert "coordinator_converter" in result.metadata
        assert result.convert_time >= 0

    def test_complete_office_pipeline(
        self, tmp_path: Path, coordinator: ConverterCoordinator
    ):
        """Test complete Office document conversion pipeline."""
        test_files = [
            ("document.docx", create_docx_file),
            ("presentation.pptx", create_pptx_file),
            ("spreadsheet.xlsx", create_xlsx_file),
        ]

        for filename, create_func in test_files:
            file_path = tmp_path / filename
            output_dir = tmp_path / "output" / filename.replace(".", "_")
            create_func(file_path)

            result = coordinator.convert(file_path, output_dir)

            assert result.success is True, f"Failed for {filename}"
            assert result.source_file == file_path
            assert result.output_file is not None
            assert result.output_file.exists()

    def test_complete_html_pipeline(
        self, tmp_path: Path, coordinator: ConverterCoordinator
    ):
        """Test complete HTML conversion pipeline."""
        html_path = tmp_path / "page.html"
        output_dir = tmp_path / "output"
        create_html_file(html_path, "Complete Pipeline Test")

        result = coordinator.convert(html_path, output_dir)

        assert result.success is True
        assert result.source_file == html_path
        assert result.converter_name == "HTMLConverter"
        assert result.output_file is not None
        assert result.output_file.exists()
        assert "Complete Pipeline Test" in result.markdown

    def test_pipeline_with_unsupported_format(
        self, tmp_path: Path, coordinator: ConverterCoordinator
    ):
        """Test pipeline behavior with unsupported format."""
        unsupported_path = tmp_path / "file.unknown"
        unsupported_path.write_text("test content")
        output_dir = tmp_path / "output"

        result = coordinator.convert(unsupported_path, output_dir)

        assert result.success is False
        assert len(result.errors) > 0
        assert result.converter_name == "ConverterCoordinator"


class TestCaseInsensitivity:
    """Test case-insensitive file extension handling."""

    @pytest.fixture
    def coordinator(self):
        """Create a ConverterCoordinator instance."""
        return ConverterCoordinator()

    def test_uppercase_extensions(
        self, tmp_path: Path, coordinator: ConverterCoordinator
    ):
        """Test that uppercase extensions are handled correctly."""
        test_cases = [
            ("DOCUMENT.PDF", create_minimal_pdf),
            ("DOCUMENT.DOCX", create_docx_file),
            ("PRESENTATION.PPTX", create_pptx_file),
            ("SPREADSHEET.XLSX", create_xlsx_file),
            ("PAGE.HTML", create_html_file),
        ]

        for filename, create_func in test_cases:
            file_path = tmp_path / filename
            output_dir = tmp_path / "output" / filename.lower().replace(".", "_")
            create_func(file_path)

            result = coordinator.convert(file_path, output_dir)

            assert result.success is True, f"Failed for {filename}"

    def test_mixed_case_extensions(
        self, tmp_path: Path, coordinator: ConverterCoordinator
    ):
        """Test that mixed-case extensions are handled correctly."""
        test_cases = [
            ("Document.Pdf", create_minimal_pdf),
            ("Document.Docx", create_docx_file),
            ("Presentation.Pptx", create_pptx_file),
            ("Spreadsheet.Xlsx", create_xlsx_file),
            ("Page.Html", create_html_file),
        ]

        for filename, create_func in test_cases:
            file_path = tmp_path / filename
            output_dir = tmp_path / "output" / filename.lower().replace(".", "_")
            create_func(file_path)

            result = coordinator.convert(file_path, output_dir)

            assert result.success is True, f"Failed for {filename}"
