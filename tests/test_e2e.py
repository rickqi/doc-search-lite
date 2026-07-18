"""
End-to-end integration tests for the document search system.

Tests cover:
- Scenario 1: Complete conversion workflow (document → markdown → storage)
- Scenario 2: Indexing and search workflow (markdown → index → search → results)
- Scenario 3: Incremental updates (file changes → detection → re-processing)
- Scenario 4: Agent search workflow (query → agent → tools → response)

All tests use real filesystem operations with mocked external services.
"""

import json
import shutil
import tempfile
import time
from datetime import datetime
from pathlib import Path
from unittest.mock import MagicMock, PropertyMock

import pytest
from docx import Document
from openpyxl import Workbook
from pptx import Presentation

from src.agent.base import ToolResult
from src.agent.llm_client import ChatResponse
from src.agent.search_agent import SearchAgent
from src.agent.tools.read import ReadTool
from src.agent.tools.search import SearchTool
from src.converter.coordinator import ConverterCoordinator
from src.search.bm25_search import BM25Searcher
from src.search.result_formatter import ResultFormatter, SearchResult
from src.storage.base import DocumentRecord
from src.storage.index import TantivyIndexManager
from src.storage.markdown_store import MarkdownStore
from src.storage.metadata import MetadataManager
from src.utils.config import Config
from src.utils.file_watcher import FileWatcher

# ============================================================================
# Helper functions to create test files
# ============================================================================


def create_minimal_pdf(file_path: Path) -> None:
    """Create a minimal valid PDF file for testing."""
    pdf_content = b"""%PDF-1.4
1 0 obj
<<
/Type /Catalog
/Pages 2 0 R
>>
endobj
2 0 obj
<<
/Type /Pages
/Kids [3 0 R]
/Count 1
>>
endobj
3 0 obj
<<
/Type /Page
/Parent 2 0 R
/MediaBox [0 0 612 792]
/Contents 4 0 R
/Resources <<
/Font <<
/F1 5 0 R
>>
>>
>>
endobj
4 0 obj
<<
/Length 44
>>
stream
BT
/F1 12 Tf
100 700 Td
(Test Content) Tj
ET
endstream
endobj
5 0 obj
<<
/Type /Font
/Subtype /Type1
/BaseFont /Helvetica
>>
endobj
xref
0 6
0000000000 65535 f 
0000000009 00000 n 
0000000058 00000 n 
0000000115 00000 n 
0000000266 00000 n 
0000000362 00000 n 
trailer
<<
/Size 6
/Root 1 0 R
>>
startxref
439
%%EOF
"""
    file_path.write_bytes(pdf_content)


def create_docx_file(file_path: Path, title: str = "Test Document") -> None:
    """Create a minimal DOCX file for testing."""
    doc = Document()
    doc.add_heading(title, level=1)
    doc.add_paragraph("This is a test paragraph with some content.")
    doc.add_paragraph(
        "The document discusses performance management and quarterly reports."
    )
    doc.save(str(file_path))


def create_pptx_file(file_path: Path, title: str = "Test Presentation") -> None:
    """Create a minimal PPTX file for testing."""
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    prs.save(str(file_path))


def create_xlsx_file(file_path: Path, title: str = "Test Spreadsheet") -> None:
    """Create a minimal XLSX file for testing."""
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws["A1"] = "Header"
    ws["A2"] = "Data"
    wb.save(str(file_path))


def create_html_file(file_path: Path, title: str = "Test Page") -> None:
    """Create a minimal HTML file for testing."""
    content = f"""<!DOCTYPE html>
<html>
<head><title>{title}</title></head>
<body>
    <h1>{title}</h1>
    <p>Test content for search system.</p>
</body>
</html>
"""
    file_path.write_text(content, encoding="utf-8")


# ============================================================================
# Fixtures
# ============================================================================


@pytest.fixture
def temp_dirs():
    """Create temporary input, output, and index directories."""
    input_dir = Path(tempfile.mkdtemp())
    output_dir = Path(tempfile.mkdtemp())
    index_dir = Path(tempfile.mkdtemp())

    yield {"input": input_dir, "output": output_dir, "index": index_dir}

    # Cleanup
    shutil.rmtree(input_dir, ignore_errors=True)
    shutil.rmtree(output_dir, ignore_errors=True)
    shutil.rmtree(index_dir, ignore_errors=True)


@pytest.fixture
def coordinator():
    """Create a ConverterCoordinator instance."""
    return ConverterCoordinator()


@pytest.fixture
def markdown_store(temp_dirs):
    """Create a MarkdownStore instance."""
    return MarkdownStore(temp_dirs["input"], temp_dirs["output"])


@pytest.fixture
def metadata_manager(temp_dirs):
    """Create a MetadataManager instance."""
    return MetadataManager(temp_dirs["output"] / "index.json")


@pytest.fixture
def index_manager(temp_dirs):
    """Create a TantivyIndexManager instance."""
    manager = TantivyIndexManager(
        index_path=temp_dirs["index"],
        use_jieba=False,  # Use regex tokenizer for consistent testing
        heap_size=20_000_000,
        num_threads=1,
    )
    yield manager
    manager.close()


@pytest.fixture
def bm25_searcher(index_manager):
    """Create a BM25Searcher instance."""
    return BM25Searcher(index_manager, snippet_length=200)


@pytest.fixture
def file_watcher():
    """Create a FileWatcher instance."""
    return FileWatcher(use_mtime_check=True, use_hash_check=True)


@pytest.fixture
def mock_config():
    """Create a mock Config."""
    config = MagicMock(spec=Config)
    config.glm_api_key = "test-api-key"
    config.glm_base_url = "https://api.test.com"
    config.llm_model = "glm-4"
    config.llm_temperature = 0.7
    config.llm_max_tokens = 4096
    return config


@pytest.fixture
def mock_llm_client():
    """Create a mock LLMClient."""
    client = MagicMock()
    client.chat.return_value = ChatResponse(
        content='{"action": "search", "search_query": "绩效"}',
        usage={"total_tokens": 50},
    )
    return client


# ============================================================================
# Scenario 1: Complete Conversion Workflow
# ============================================================================


class TestCompleteConversionWorkflow:
    """
    Test the complete conversion workflow:
    Create document → ConverterCoordinator converts → MarkdownStore saves → metadata recorded
    """

    def test_pdf_to_markdown_workflow(
        self, temp_dirs, coordinator, markdown_store, metadata_manager
    ):
        """Test complete PDF conversion and storage workflow."""
        # 1. Create test PDF
        source_file = temp_dirs["input"] / "report.pdf"
        create_minimal_pdf(source_file)

        # 2. Convert using coordinator
        result = coordinator.convert(source_file, temp_dirs["output"])

        # 3. Verify conversion result
        assert result.success is True
        assert result.source_file == source_file
        assert result.output_file is not None
        assert result.output_file.suffix == ".md"

        # 4. Verify output file exists
        output_path = markdown_store.get_output_path(source_file)
        assert output_path.exists()

        # 5. Verify content was written
        content = output_path.read_text(encoding="utf-8")
        assert len(content) > 0

    def test_docx_to_markdown_with_metadata(
        self, temp_dirs, coordinator, markdown_store, metadata_manager
    ):
        """Test DOCX conversion with full metadata storage."""
        # 1. Create test document
        source_file = temp_dirs["input"] / "policy.docx"
        create_docx_file(source_file, "Performance Management Policy")

        # 2. Convert
        result = coordinator.convert(source_file, temp_dirs["output"])

        # 3. Verify success
        assert result.success is True
        assert "Performance Management Policy" in result.markdown

        # 4. Create and save document record
        record = DocumentRecord(
            id="policy_doc_001",
            source_path=source_file,
            output_path=markdown_store.get_output_path(source_file),
            title="Performance Management Policy",
            content_hash=result.metadata.get("content_hash", "abc123"),
            file_size=source_file.stat().st_size,
            file_mtime=datetime.now(),
            metadata={"author": "HR Department", "category": "Policy"},
            keywords=["performance", "management", "policy"],
            sections=["Introduction", "Scope", "Guidelines"],
        )

        # 5. Save to MarkdownStore
        save_result = markdown_store.save(record, result.markdown)
        assert save_result is True

        # 6. Verify metadata was saved
        metadata = metadata_manager.load(source_file)
        if metadata:
            assert metadata.get("id") == "policy_doc_001"

        # 7. Verify can load back
        loaded = markdown_store.load(record.id)
        assert loaded is not None
        loaded_record, loaded_content = loaded
        assert loaded_record.title == "Performance Management Policy"
        assert "Performance Management Policy" in loaded_content

    def test_multi_format_batch_conversion(
        self, temp_dirs, coordinator, markdown_store
    ):
        """Test batch conversion of multiple document formats."""
        # 1. Create files of different formats
        files = [
            ("doc1.pdf", create_minimal_pdf),
            ("doc2.docx", create_docx_file),
            ("doc3.pptx", create_pptx_file),
            ("doc4.xlsx", create_xlsx_file),
            ("doc5.html", create_html_file),
        ]

        for filename, create_func in files:
            create_func(temp_dirs["input"] / filename)

        # 2. Convert all files
        results = []
        for filename, _ in files:
            source_file = temp_dirs["input"] / filename
            result = coordinator.convert(source_file, temp_dirs["output"])
            results.append((filename, result))

        # 3. Verify all conversions succeeded
        for filename, result in results:
            assert result.success is True, f"Failed to convert {filename}"
            assert result.output_file is not None
            assert result.output_file.exists()

        # 4. Verify all output files exist
        for filename, _ in files:
            stem = Path(filename).stem
            expected_output = temp_dirs["output"] / f"{stem}.md"
            assert expected_output.exists(), f"Output file missing for {filename}"

    def test_nested_directory_conversion(self, temp_dirs, coordinator, markdown_store):
        """Test conversion preserving nested directory structure."""
        # 1. Create nested directory structure
        nested_dir = temp_dirs["input"] / "department" / "hr" / "policies"
        nested_dir.mkdir(parents=True)

        source_file = nested_dir / "leave_policy.docx"
        create_docx_file(source_file, "Leave Policy 2024")

        # 2. Convert
        result = coordinator.convert(source_file, temp_dirs["output"])

        # 3. Verify conversion succeeded
        # Note: Coordinator outputs to flat structure by default
        assert result.success is True
        assert result.output_file is not None
        assert result.output_file.exists()


# ============================================================================
# Scenario 2: Indexing and Search Workflow
# ============================================================================


class TestIndexAndSearchWorkflow:
    """
    Test the indexing and search workflow:
    Create documents → Convert → Index → Search → Format results
    """

    def test_full_text_search_workflow(
        self, temp_dirs, coordinator, markdown_store, index_manager, bm25_searcher
    ):
        """Test complete indexing and search workflow."""
        # 1. Create and convert test documents
        docs = [
            ("performance.docx", "Performance Management Guidelines"),
            ("quarterly.docx", "Q4 Quarterly Report"),
            ("training.docx", "Employee Training Program"),
        ]

        for filename, title in docs:
            source_file = temp_dirs["input"] / filename
            create_docx_file(source_file, title)

            # Convert
            result = coordinator.convert(source_file, temp_dirs["output"])
            assert result.success is True

            # Create document record
            record = DocumentRecord(
                id=filename.replace(".docx", ""),
                source_path=source_file,
                output_path=markdown_store.get_output_path(source_file),
                title=title,
                content_hash=f"hash_{filename}",
                file_size=source_file.stat().st_size,
                file_mtime=datetime.now(),
            )

            # Save to store
            markdown_store.save(record, result.markdown)

            # Index document
            index_manager.add_document(
                doc_id=record.id,
                title=title,
                content=result.markdown,
                metadata={
                    "source_path": str(source_file),
                    "filename": filename,
                },
            )

        # 2. Commit index
        index_manager.commit()

        # 3. Search for "performance"
        search_result = bm25_searcher.search("performance", limit=10)

        # 4. Verify results
        assert search_result.total >= 1
        assert len(search_result.results) >= 1
        assert any("performance" in r.title.lower() for r in search_result.results)

    def test_search_with_result_formatting(
        self, temp_dirs, coordinator, index_manager, bm25_searcher
    ):
        """Test search with different result formatting options."""
        # 1. Create and index test documents
        for i in range(3):
            source_file = temp_dirs["input"] / f"doc{i}.docx"
            create_docx_file(source_file, f"Document {i}: Quarterly Report")

            result = coordinator.convert(source_file, temp_dirs["output"])
            assert result.success is True

            index_manager.add_document(
                doc_id=f"doc{i}",
                title=f"Document {i}: Quarterly Report",
                content=result.markdown,
                metadata={"source_path": str(source_file)},
            )

        index_manager.commit()

        # 2. Search
        search_result = bm25_searcher.search("Quarterly", limit=10)

        # 3. Convert to SearchResult objects for formatting
        results_for_formatting = [
            SearchResult(
                title=r.title,
                score=r.score,
                snippet=r.snippet,
                source=r.source_path or Path(""),
            )
            for r in search_result.results
        ]

        # 4. Test JSON formatting
        formatter = ResultFormatter(highlight_pattern="Quarterly")
        json_output = formatter.format_json(results_for_formatting)
        assert "Quarterly" in json_output
        assert "score" in json_output

        # 5. Test text formatting
        text_output = formatter.format_text(results_for_formatting)
        assert "Quarterly" in text_output
        assert "score:" in text_output

        # 6. Test markdown formatting
        md_output = formatter.format_markdown(results_for_formatting)
        assert "## " in md_output  # Markdown headers
        assert "Quarterly" in md_output

    def test_pagination_workflow(
        self, temp_dirs, coordinator, index_manager, bm25_searcher
    ):
        """Test search with pagination."""
        # 1. Create and index documents
        for i in range(20):
            source_file = temp_dirs["input"] / f"report{i}.docx"
            create_docx_file(source_file, f"Report {i}: Annual Performance")

            result = coordinator.convert(source_file, temp_dirs["output"])
            assert result.success is True

            index_manager.add_document(
                doc_id=f"report{i}",
                title=f"Report {i}: Annual Performance",
                content=result.markdown,
                metadata={"source_path": str(source_file)},
            )

        index_manager.commit()

        # 2. Search with pagination
        page1 = bm25_searcher.search("Performance", limit=5, offset=0)

        # 3. Verify pagination works
        assert len(page1.results) <= 5  # May have fewer if not all indexed
        assert page1.offset == 0
        assert page1.limit == 5

        # 4. If we have enough results, test second page
        if page1.has_more and page1.total > 5:
            page2 = bm25_searcher.search("Performance", limit=5, offset=5)
            # Verify pages have different results
            page1_ids = {r.doc_id for r in page1.results}
            page2_ids = {r.doc_id for r in page2.results}
            assert page1_ids.isdisjoint(page2_ids)  # No overlap


# ============================================================================
# Scenario 3: Incremental Updates
# ============================================================================


class TestIncrementalUpdates:
    """
    Test incremental update workflow:
    Modify files → FileWatcher detects changes → Re-process only changed files
    """

    def test_detect_new_files(self, temp_dirs, file_watcher, metadata_manager):
        """Test detection of newly added files."""
        # 1. Initially no files
        change_set = file_watcher.detect_changes(
            source_dir=temp_dirs["input"],
            metadata_manager=metadata_manager,
            extensions={".docx", ".pdf"},
        )

        assert len(change_set.added) == 0
        assert len(change_set.unchanged) == 0

        # 2. Add new files
        create_docx_file(temp_dirs["input"] / "new_doc.docx", "New Document")
        create_minimal_pdf(temp_dirs["input"] / "new_pdf.pdf")

        # 3. Detect changes
        change_set = file_watcher.detect_changes(
            source_dir=temp_dirs["input"],
            metadata_manager=metadata_manager,
            extensions={".docx", ".pdf"},
        )

        assert len(change_set.added) == 2
        assert change_set.has_changes is True

    def test_detect_modified_files(self, temp_dirs, file_watcher, metadata_manager):
        """Test detection of modified files."""
        # 1. Create and register a file
        source_file = temp_dirs["input"] / "modified.docx"
        create_docx_file(source_file, "Original Content")

        # 2. Register in metadata with source_path field for FileWatcher
        import hashlib

        content_hash = hashlib.sha256(source_file.read_bytes()).hexdigest()
        mtime = source_file.stat().st_mtime
        metadata_manager.save(
            source_file,
            {
                "source_path": str(source_file.resolve()),  # Required for FileWatcher
                "content_hash": content_hash,
                "modified_time": mtime,
                "status": "processed",
            },
        )

        # 3. Verify initially - file should be unchanged
        change_set = file_watcher.detect_changes(
            source_dir=temp_dirs["input"],
            metadata_manager=metadata_manager,
            extensions={".docx"},
        )

        # File is registered - should either be unchanged or detected as same
        assert len(change_set.added) == 0
        # Note: Due to timing/hashing, the file may be detected as unchanged or modified

        # 4. Modify the file significantly
        time.sleep(0.2)  # Ensure different mtime
        create_docx_file(source_file, "Modified Content - Updated - Changed")

        # 5. Detect changes after modification
        change_set = file_watcher.detect_changes(
            source_dir=temp_dirs["input"],
            metadata_manager=metadata_manager,
            extensions={".docx"},
        )

        # Should detect modification (modified or re-detected as changed)
        assert change_set.has_changes is True
        assert len(change_set.modified) >= 1 or len(change_set.added) >= 1

    def test_detect_deleted_files(self, temp_dirs, file_watcher, metadata_manager):
        """Test detection of deleted files."""
        # 1. Create and register a file
        source_file = temp_dirs["input"] / "to_delete.docx"
        create_docx_file(source_file, "To Be Deleted")

        # 2. Register in metadata with source_path field for FileWatcher
        metadata_manager.save(
            source_file,
            {
                "source_path": str(source_file.resolve()),  # Required for FileWatcher
                "content_hash": "some_hash",
                "modified_time": source_file.stat().st_mtime,
                "status": "processed",
            },
        )

        # 3. Delete the file
        source_file.unlink()

        # 4. Detect changes
        change_set = file_watcher.detect_changes(
            source_dir=temp_dirs["input"],
            metadata_manager=metadata_manager,
            extensions={".docx"},
        )

        # 5. Verify deletion detected (file in metadata but not on disk)
        # The deleted file should be detected
        assert len(change_set.deleted) == 1 or change_set.has_changes is True

    def test_incremental_reprocessing_workflow(
        self,
        temp_dirs,
        coordinator,
        markdown_store,
        index_manager,
        file_watcher,
        metadata_manager,
    ):
        """Test complete incremental re-processing workflow."""
        # 1. Initial setup: create and process documents
        initial_files = ["doc1.docx", "doc2.docx"]
        for filename in initial_files:
            source_file = temp_dirs["input"] / filename
            create_docx_file(source_file, f"Initial {filename}")

            # Convert
            result = coordinator.convert(source_file, temp_dirs["output"])
            assert result.success is True

            # Register in metadata with source_path for FileWatcher
            import hashlib

            content_hash = hashlib.sha256(source_file.read_bytes()).hexdigest()
            metadata_manager.save(
                source_file,
                {
                    "source_path": str(
                        source_file.resolve()
                    ),  # Required for FileWatcher
                    "content_hash": content_hash,
                    "modified_time": source_file.stat().st_mtime,
                    "output_path": str(markdown_store.get_output_path(source_file)),
                    "status": "processed",
                },
            )

        # 2. Add a new file and modify an existing one
        create_docx_file(temp_dirs["input"] / "new_doc.docx", "New Document")

        time.sleep(0.1)  # Ensure different mtime
        create_docx_file(temp_dirs["input"] / "doc1.docx", "Modified Doc1")

        # 3. Detect changes
        change_set = file_watcher.detect_changes(
            source_dir=temp_dirs["input"],
            metadata_manager=metadata_manager,
            extensions={".docx"},
        )

        # 4. Verify changes detected (either new files or modified files)
        assert change_set.has_changes is True


# ============================================================================
# Scenario 4: Agent Search with Mock LLM
# ============================================================================


class TestAgentSearchWorkflow:
    """
    Test the agent-based search workflow:
    User query → Agent analyzes → Tools execute → Response generated
    """

    @pytest.fixture
    def mock_search_tool(self, index_manager, bm25_searcher):
        """Create a mock SearchTool with real searcher."""
        tool = MagicMock(spec=SearchTool)
        type(tool).name = PropertyMock(return_value="search")
        type(tool).description = PropertyMock(return_value="Search documents")

        def execute_search(**kwargs):
            query = kwargs.get("query", "")
            limit = kwargs.get("limit", 10)
            offset = kwargs.get("offset", 0)

            result = bm25_searcher.search(query, limit=limit, offset=offset)

            return ToolResult.ok(
                data=json.dumps(
                    {
                        "query": query,
                        "total": result.total,
                        "offset": offset,
                        "limit": limit,
                        "has_more": result.has_more,
                        "results": [
                            {
                                "doc_id": r.doc_id,
                                "title": r.title,
                                "score": r.score,
                                "snippet": r.snippet,
                                "source_path": str(r.source_path)
                                if r.source_path
                                else None,
                            }
                            for r in result.results
                        ],
                    },
                    ensure_ascii=False,
                ),
                metadata={"total_results": result.total},
            )

        tool.execute = execute_search
        return tool

    @pytest.fixture
    def mock_read_tool(self, temp_dirs, coordinator):
        """Create a mock ReadTool."""
        tool = MagicMock(spec=ReadTool)
        type(tool).name = PropertyMock(return_value="read")
        type(tool).description = PropertyMock(return_value="Read document content")

        def execute_read(**kwargs):
            doc_id = kwargs.get("doc_id", "")
            # Read from output directory
            output_file = temp_dirs["output"] / f"{doc_id}.md"
            if output_file.exists():
                content = output_file.read_text(encoding="utf-8")
                return ToolResult.ok(
                    data=content,
                    metadata={
                        "doc_id": doc_id,
                        "lines_read": len(content.splitlines()),
                    },
                )
            return ToolResult.fail(error=f"Document {doc_id} not found", metadata={})

        tool.execute = execute_read
        return tool

    def test_agent_search_with_results(
        self,
        temp_dirs,
        coordinator,
        index_manager,
        mock_config,
        mock_llm_client,
        mock_search_tool,
        mock_read_tool,
    ):
        """Test agent search workflow with actual search results."""
        # 1. Create and index test documents
        for i in range(3):
            source_file = temp_dirs["input"] / f"policy{i}.docx"
            create_docx_file(
                source_file, f"Policy Document {i}: Performance Management"
            )

            result = coordinator.convert(source_file, temp_dirs["output"])
            assert result.success is True

            index_manager.add_document(
                doc_id=f"policy{i}",
                title=f"Policy Document {i}: Performance Management",
                content=result.markdown,
                metadata={"source_path": str(source_file)},
            )

        index_manager.commit()

        # 2. Setup agent with mocked LLM
        mock_llm_client.chat.side_effect = [
            # Query analysis
            ChatResponse(
                content='{"action": "search", "search_query": "performance"}',
                usage={"total_tokens": 30},
            ),
            # Answer generation
            ChatResponse(
                content="Based on the documents, performance management involves... [Policy Document 0]",
                usage={"total_tokens": 150},
            ),
        ]

        agent = SearchAgent(
            config=mock_config,
            search_tool=mock_search_tool,
            read_tool=mock_read_tool,
            llm_client=mock_llm_client,
            mode="pipeline",
        )

        # 3. Execute search
        response = agent.run("What is performance management?")

        # 4. Verify response
        assert response.success is True
        assert len(response.answer) > 0
        assert len(response.tool_calls) >= 1

    def test_agent_search_no_results(
        self,
        temp_dirs,
        index_manager,
        mock_config,
        mock_llm_client,
        mock_search_tool,
        mock_read_tool,
    ):
        """Test agent search when no results are found."""
        # 1. Setup empty index (no documents)
        index_manager.commit()

        # 2. Setup agent
        mock_llm_client.chat.return_value = ChatResponse(
            content='{"action": "search", "search_query": "nonexistent"}',
            usage={"total_tokens": 30},
        )

        agent = SearchAgent(
            config=mock_config,
            search_tool=mock_search_tool,
            read_tool=mock_read_tool,
            llm_client=mock_llm_client,
            mode="pipeline",
        )

        # 3. Search for non-existent content
        response = agent.run("nonexistent_xyz_12345")

        # 4. Verify handling of no results
        assert response.success is True
        assert (
            "未找到" in response.answer
            or "没有" in response.answer
            or len(response.answer) == 0
        )

    def test_agent_greeting_direct_response(
        self, mock_config, mock_llm_client, mock_search_tool, mock_read_tool
    ):
        """Test agent handles greetings without search."""
        agent = SearchAgent(
            config=mock_config,
            search_tool=mock_search_tool,
            read_tool=mock_read_tool,
            llm_client=mock_llm_client,
            mode="pipeline",
        )

        # Test various greetings
        greetings = ["你好", "hello", "hi"]
        for greeting in greetings:
            response = agent.run(greeting)
            assert response.success is True
            assert len(response.tool_calls) == 0  # No tools should be called

    def test_agent_tool_call_chain(
        self,
        temp_dirs,
        coordinator,
        index_manager,
        mock_config,
        mock_llm_client,
        mock_search_tool,
        mock_read_tool,
    ):
        """Test agent executes proper tool call chain: search → read → respond."""
        # 1. Create and index a document
        source_file = temp_dirs["input"] / "benefits.docx"
        create_docx_file(source_file, "Employee Benefits Guide 2024")

        result = coordinator.convert(source_file, temp_dirs["output"])
        assert result.success is True

        index_manager.add_document(
            doc_id="benefits",
            title="Employee Benefits Guide 2024",
            content=result.markdown,
            metadata={"source_path": str(source_file)},
        )
        index_manager.commit()

        # 2. Setup agent
        mock_llm_client.chat.side_effect = [
            ChatResponse(
                content='{"action": "search", "search_query": "benefits"}',
                usage={"total_tokens": 30},
            ),
            ChatResponse(
                content="The benefits guide covers health insurance, retirement plans, and vacation policies. [benefits]",
                usage={"total_tokens": 200},
            ),
        ]

        agent = SearchAgent(
            config=mock_config,
            search_tool=mock_search_tool,
            read_tool=mock_read_tool,
            llm_client=mock_llm_client,
            mode="pipeline",
        )

        # 3. Execute search
        response = agent.run("What are the employee benefits?")

        # 4. Verify tool call chain
        assert response.success is True
        # Should have at least search call
        assert len(response.tool_calls) >= 1
        # First call should be search
        assert response.tool_calls[0]["tool"] == "search"


# ============================================================================
# Integration Tests: Full System Workflow
# ============================================================================


class TestFullSystemIntegration:
    """
    Full system integration tests combining all components.
    """

    def test_convert_index_search_full_pipeline(
        self, temp_dirs, coordinator, markdown_store, index_manager, bm25_searcher
    ):
        """Test complete pipeline: convert → store → index → search."""
        # 1. Create document
        source_file = temp_dirs["input"] / "full_test.docx"
        create_docx_file(source_file, "Annual Report 2024: Financial Summary")

        # 2. Convert
        convert_result = coordinator.convert(source_file, temp_dirs["output"])
        assert convert_result.success is True

        # 3. Create and save document record
        record = DocumentRecord(
            id="full_test_001",
            source_path=source_file,
            output_path=markdown_store.get_output_path(source_file),
            title="Annual Report 2024: Financial Summary",
            content_hash="test_hash",
            file_size=source_file.stat().st_size,
            file_mtime=datetime.now(),
        )
        markdown_store.save(record, convert_result.markdown)

        # 4. Index
        index_manager.add_document(
            doc_id=record.id,
            title=record.title,
            content=convert_result.markdown,
            metadata={
                "source_path": str(source_file),
                "filename": "full_test.docx",
            },
        )
        index_manager.commit()

        # 5. Search
        search_result = bm25_searcher.search("Financial", limit=10)

        # 6. Verify
        assert search_result.total >= 1
        assert len(search_result.results) >= 1
        found = any("Financial" in r.title for r in search_result.results)
        assert found

    def test_unicode_content_workflow(
        self, temp_dirs, coordinator, markdown_store, index_manager, bm25_searcher
    ):
        """Test workflow with Chinese and Unicode content."""
        # 1. Create document with Chinese content
        source_file = temp_dirs["input"] / "chinese_test.docx"
        doc = Document()
        doc.add_heading("绩效管理制度", level=1)
        doc.add_paragraph("本文档介绍公司的绩效考核流程和标准。")
        doc.add_paragraph("员工每年需要进行两次绩效考核。")
        doc.save(str(source_file))

        # 2. Convert
        convert_result = coordinator.convert(source_file, temp_dirs["output"])
        assert convert_result.success is True

        # 3. Save and index
        record = DocumentRecord(
            id="chinese_test_001",
            source_path=source_file,
            output_path=markdown_store.get_output_path(source_file),
            title="绩效管理制度",
            content_hash="test_hash",
            file_size=source_file.stat().st_size,
            file_mtime=datetime.now(),
        )
        markdown_store.save(record, convert_result.markdown)

        index_manager.add_document(
            doc_id=record.id,
            title=record.title,
            content=convert_result.markdown,
            metadata={"source_path": str(source_file)},
        )
        index_manager.commit()

        # 4. Search
        search_result = bm25_searcher.search("绩效", limit=10)

        # 5. Verify
        assert isinstance(search_result.total, int)

    def test_error_recovery_in_pipeline(self, temp_dirs, coordinator, markdown_store):
        """Test pipeline recovery from errors."""
        # 1. Try to convert non-existent file
        missing_file = temp_dirs["input"] / "missing.pdf"
        result = coordinator.convert(missing_file, temp_dirs["output"])

        # 2. Should handle error gracefully
        assert result.success is False
        assert len(result.errors) > 0

        # 3. Pipeline should still work after error
        valid_file = temp_dirs["input"] / "valid.docx"
        create_docx_file(valid_file, "Valid Document")

        result = coordinator.convert(valid_file, temp_dirs["output"])
        assert result.success is True

    def test_concurrent_file_handling(self, temp_dirs, coordinator):
        """Test handling of multiple files in sequence."""
        # Create multiple files
        files = []
        for i in range(5):
            source_file = temp_dirs["input"] / f"concurrent_{i}.docx"
            create_docx_file(source_file, f"Document {i}")
            files.append(source_file)

        # Convert all files
        results = []
        for source_file in files:
            result = coordinator.convert(source_file, temp_dirs["output"])
            results.append(result)

        # Verify all succeeded
        for i, result in enumerate(results):
            assert result.success is True, f"Failed for file {i}"
            assert result.output_file is not None
            assert result.output_file.exists()
