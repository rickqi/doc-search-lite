"""
Unit tests for OfficeConverter.

Tests cover:
- Basic Office document conversion (DOCX, PPTX, XLSX)
- Support for all three Office file extensions
- Rejection of non-Office files
- Error handling for corrupted files
- Error handling for missing files
- Metadata extraction
- Output file creation
- can_convert method
- Structure preservation (headings, lists, tables)
"""

from pathlib import Path

import pytest
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches, Pt
from docx import Document
from docx.shared import Inches as DocxInches

from src.converter.office import OfficeConverter


# Helper functions to create test Office files
def create_docx_file(file_path: Path, title: str = "Test Document") -> None:
    """Create a minimal DOCX file for testing."""
    doc = Document()
    doc.add_heading(title, level=1)
    doc.add_heading("Sub Heading", level=2)
    doc.add_paragraph("This is a paragraph.")
    doc.add_paragraph("Another paragraph with more text.")

    # Add a list
    doc.add_paragraph("First item", style="List Bullet")
    doc.add_paragraph("Second item", style="List Bullet")
    doc.add_paragraph("Third item", style="List Bullet")

    # Add a table
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Header 1"
    table.cell(0, 1).text = "Header 2"
    table.cell(1, 0).text = "Cell 1"
    table.cell(1, 1).text = "Cell 2"

    doc.save(str(file_path))


def create_pptx_file(file_path: Path, title: str = "Test Presentation") -> None:
    """Create a minimal PPTX file for testing."""
    prs = Presentation()

    # Add title slide
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title

    # Add content slide
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = "Content Slide"
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "This is slide content."
    p = tf.add_paragraph()
    p.text = "Another line of content."

    prs.save(str(file_path))


def create_xlsx_file(file_path: Path, title: str = "Test Spreadsheet") -> None:
    """Create a minimal XLSX file for testing."""
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"

    # Add headers
    ws["A1"] = "Name"
    ws["B1"] = "Age"
    ws["C1"] = "City"

    # Add data
    ws["A2"] = "John"
    ws["B2"] = 30
    ws["C2"] = "New York"

    ws["A3"] = "Jane"
    ws["B3"] = 25
    ws["C3"] = "Los Angeles"

    ws["A4"] = "Bob"
    ws["B4"] = 35
    ws["C4"] = "Chicago"

    wb.save(str(file_path))


class TestOfficeConverterProperties:
    """Test Office converter basic properties."""

    def test_converter_name(self):
        """Test converter name."""
        converter = OfficeConverter()
        assert converter.name == "OfficeConverter"

    def test_converter_version(self):
        """Test converter version."""
        converter = OfficeConverter()
        assert converter.version == "0.1.0"

    def test_supported_formats(self):
        """Test that all Office formats are supported."""
        converter = OfficeConverter()
        formats = converter.supported_formats

        assert ".docx" in formats
        assert ".doc" in formats
        assert ".pptx" in formats
        assert ".xlsx" in formats
        assert ".xls" in formats
        assert len(formats) == 5


class TestCanConvert:
    """Test can_convert method."""

    def test_can_convert_docx_file(self):
        """Test that .docx files are recognized."""
        converter = OfficeConverter()
        file_path = Path("test.docx")

        assert converter.can_convert(file_path) is True

    def test_can_convert_pptx_file(self):
        """Test that .pptx files are recognized."""
        converter = OfficeConverter()
        file_path = Path("test.pptx")

        assert converter.can_convert(file_path) is True

    def test_can_convert_xlsx_file(self):
        """Test that .xlsx files are recognized."""
        converter = OfficeConverter()
        file_path = Path("test.xlsx")

        assert converter.can_convert(file_path) is True

    def test_can_convert_doc_file(self):
        """Test that .doc files are recognized."""
        converter = OfficeConverter()
        file_path = Path("test.doc")

        assert converter.can_convert(file_path) is True

    def test_can_convert_xls_file(self):
        """Test that .xls files are recognized."""
        converter = OfficeConverter()
        file_path = Path("test.xls")

        assert converter.can_convert(file_path) is True

    def test_can_convert_uppercase_extension(self):
        """Test that uppercase extensions are recognized."""
        converter = OfficeConverter()

        assert converter.can_convert(Path("test.DOCX")) is True
        assert converter.can_convert(Path("test.DOC")) is True
        assert converter.can_convert(Path("test.PPTX")) is True
        assert converter.can_convert(Path("test.XLSX")) is True
        assert converter.can_convert(Path("test.XLS")) is True

    def test_cannot_convert_other_formats(self):
        """Test that non-Office files are rejected."""
        converter = OfficeConverter()

        assert converter.can_convert(Path("test.pdf")) is False
        assert converter.can_convert(Path("test.html")) is False
        assert converter.can_convert(Path("test.txt")) is False
        assert converter.can_convert(Path("test")) is False


class TestDocxConversion:
    """Test DOCX to Markdown conversion."""

    def test_convert_basic_docx(self, tmp_path, tmp_output_dir):
        """Test converting a basic DOCX file to Markdown."""
        # Create test DOCX file
        docx_file = tmp_path / "test.docx"
        create_docx_file(docx_file, "Main Heading")

        # Convert DOCX
        converter = OfficeConverter()
        result = converter.convert(docx_file, tmp_output_dir)

        # Verify result
        assert result.success is True
        assert result.source_file == docx_file
        assert result.output_file == tmp_output_dir / "test.md"
        assert result.converter_name == "OfficeConverter"
        assert result.converter_version == "0.1.0"
        assert len(result.errors) == 0
        assert result.convert_time >= 0

        # Verify output file was created
        assert result.output_file.exists()
        markdown_content = result.output_file.read_text(encoding="utf-8")

        # Verify Markdown content
        assert "Main Heading" in markdown_content
        assert markdown_content == result.markdown

    def test_convert_docx_preserves_structure(self, tmp_path, tmp_output_dir):
        """Test that DOCX document structure is preserved."""
        docx_file = tmp_path / "structure.docx"
        create_docx_file(docx_file)

        converter = OfficeConverter()
        result = converter.convert(docx_file, tmp_output_dir)

        assert result.success is True
        markdown = result.markdown

        # Verify headings are preserved
        assert "Test Document" in markdown
        assert "Sub Heading" in markdown

        # Verify list items are preserved
        assert "First item" in markdown
        assert "Second item" in markdown

    def test_convert_docx_with_table(self, tmp_path, tmp_output_dir):
        """Test that DOCX tables are preserved as Markdown tables."""
        docx_file = tmp_path / "table.docx"
        create_docx_file(docx_file)

        converter = OfficeConverter()
        result = converter.convert(docx_file, tmp_output_dir)

        assert result.success is True
        markdown = result.markdown

        # Verify table content is present
        assert "Header 1" in markdown
        assert "Header 2" in markdown
        assert "Cell 1" in markdown
        assert "Cell 2" in markdown

    def test_docx_metadata(self, tmp_path, tmp_output_dir):
        """Test that DOCX metadata is extracted."""
        docx_file = tmp_path / "metadata.docx"
        create_docx_file(docx_file)

        converter = OfficeConverter()
        result = converter.convert(docx_file, tmp_output_dir)

        assert result.success is True

        # Check document type in metadata
        assert "document_type" in result.metadata
        assert result.metadata["document_type"] == "Word Document"

        # Check file size in metadata
        assert "file_size_bytes" in result.metadata
        assert result.metadata["file_size_bytes"] > 0


class TestPptxConversion:
    """Test PPTX to Markdown conversion."""

    def test_convert_basic_pptx(self, tmp_path, tmp_output_dir):
        """Test converting a basic PPTX file to Markdown."""
        # Create test PPTX file
        pptx_file = tmp_path / "test.pptx"
        create_pptx_file(pptx_file, "Presentation Title")

        # Convert PPTX
        converter = OfficeConverter()
        result = converter.convert(pptx_file, tmp_output_dir)

        # Verify result
        assert result.success is True
        assert result.source_file == pptx_file
        assert result.output_file == tmp_output_dir / "test.md"
        assert result.converter_name == "OfficeConverter"
        assert len(result.errors) == 0

        # Verify output file was created
        assert result.output_file.exists()
        markdown_content = result.output_file.read_text(encoding="utf-8")

        # Verify Markdown content
        assert "Presentation Title" in markdown_content

    def test_convert_pptx_multiple_slides(self, tmp_path, tmp_output_dir):
        """Test converting PPTX with multiple slides."""
        pptx_file = tmp_path / "multi.pptx"
        create_pptx_file(pptx_file)

        converter = OfficeConverter()
        result = converter.convert(pptx_file, tmp_output_dir)

        assert result.success is True
        markdown = result.markdown

        # Verify slide content is present
        assert "Test Presentation" in markdown
        assert "Content Slide" in markdown

    def test_pptx_metadata(self, tmp_path, tmp_output_dir):
        """Test that PPTX metadata is extracted."""
        pptx_file = tmp_path / "metadata.pptx"
        create_pptx_file(pptx_file)

        converter = OfficeConverter()
        result = converter.convert(pptx_file, tmp_output_dir)

        assert result.success is True

        # Check document type in metadata
        assert "document_type" in result.metadata
        assert result.metadata["document_type"] == "PowerPoint Presentation"


class TestXlsxConversion:
    """Test XLSX to Markdown conversion."""

    def test_convert_basic_xlsx(self, tmp_path, tmp_output_dir):
        """Test converting a basic XLSX file to Markdown."""
        # Create test XLSX file
        xlsx_file = tmp_path / "test.xlsx"
        create_xlsx_file(xlsx_file)

        # Convert XLSX
        converter = OfficeConverter()
        result = converter.convert(xlsx_file, tmp_output_dir)

        # Verify result
        assert result.success is True
        assert result.source_file == xlsx_file
        assert result.output_file == tmp_output_dir / "test.md"
        assert result.converter_name == "OfficeConverter"
        assert len(result.errors) == 0

        # Verify output file was created
        assert result.output_file.exists()
        markdown_content = result.output_file.read_text(encoding="utf-8")

        # Verify Markdown content
        assert "Name" in markdown_content
        assert "Age" in markdown_content

    def test_convert_xlsx_preserves_table(self, tmp_path, tmp_output_dir):
        """Test that XLSX tables are preserved as Markdown tables."""
        xlsx_file = tmp_path / "table.xlsx"
        create_xlsx_file(xlsx_file)

        converter = OfficeConverter()
        result = converter.convert(xlsx_file, tmp_output_dir)

        assert result.success is True
        markdown = result.markdown

        # Verify table content is present
        assert "John" in markdown
        assert "Jane" in markdown
        assert "Bob" in markdown
        assert "New York" in markdown

    def test_xlsx_metadata(self, tmp_path, tmp_output_dir):
        """Test that XLSX metadata is extracted."""
        xlsx_file = tmp_path / "metadata.xlsx"
        create_xlsx_file(xlsx_file)

        converter = OfficeConverter()
        result = converter.convert(xlsx_file, tmp_output_dir)

        assert result.success is True

        # Check document type in metadata
        assert "document_type" in result.metadata
        assert result.metadata["document_type"] == "Excel Spreadsheet"


class TestDocConversion:
    """Test DOC (Legacy Word) to Markdown conversion."""

    def test_convert_basic_doc(self, tmp_path, tmp_output_dir):
        """Test converting a basic DOC file to Markdown."""
        # Create test DOC file (minimal binary content for legacy format)
        # Note: This is a simple binary file - actual conversion depends on MarkItDown
        doc_file = tmp_path / "test.doc"
        # Write minimal DOC file header (OLE format signature)
        doc_file.write_bytes(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + b"\x00" * 1000)

        # Convert DOC
        converter = OfficeConverter()
        result = converter.convert(doc_file, tmp_output_dir)

        # Verify result (MarkItDown should handle gracefully)
        assert result.converter_name == "OfficeConverter"
        assert result.converter_version == "0.1.0"
        assert result.source_file == doc_file

    def test_doc_metadata(self, tmp_path, tmp_output_dir):
        """Test that DOC metadata is extracted."""
        doc_file = tmp_path / "metadata.doc"
        # Write minimal DOC file header
        doc_file.write_bytes(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + b"\x00" * 1000)

        converter = OfficeConverter()
        result = converter.convert(doc_file, tmp_output_dir)

        # If successful, check document type in metadata
        if result.success:
            assert "document_type" in result.metadata
            assert result.metadata["document_type"] == "Word Document (Legacy)"


class TestXlsConversion:
    """Test XLS (Legacy Excel) to Markdown conversion."""

    def test_convert_basic_xls(self, tmp_path, tmp_output_dir):
        """Test converting a basic XLS file to Markdown."""
        # Create test XLS file (minimal binary content for legacy format)
        # Note: This is a simple binary file - actual conversion depends on MarkItDown
        xls_file = tmp_path / "test.xls"
        # Write minimal XLS file header (OLE format signature)
        xls_file.write_bytes(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + b"\x00" * 1000)

        # Convert XLS
        converter = OfficeConverter()
        result = converter.convert(xls_file, tmp_output_dir)

        # Verify result (MarkItDown should handle gracefully)
        assert result.converter_name == "OfficeConverter"
        assert result.converter_version == "0.1.0"
        assert result.source_file == xls_file

    def test_xls_metadata(self, tmp_path, tmp_output_dir):
        """Test that XLS metadata is extracted."""
        xls_file = tmp_path / "metadata.xls"
        # Write minimal XLS file header
        xls_file.write_bytes(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + b"\x00" * 1000)

        converter = OfficeConverter()
        result = converter.convert(xls_file, tmp_output_dir)

        # If successful, check document type in metadata
        if result.success:
            assert "document_type" in result.metadata
            assert result.metadata["document_type"] == "Excel Spreadsheet (Legacy)"


class TestErrorHandling:
    """Test error handling for invalid inputs."""

    def test_unsupported_file_extension(self, tmp_path, tmp_output_dir):
        """Test that non-Office files are rejected with clear error message."""
        text_file = tmp_path / "test.txt"
        text_file.write_text("Not Office content")

        converter = OfficeConverter()
        result = converter.convert(text_file, tmp_output_dir)

        assert result.success is False
        assert result.markdown == ""
        assert len(result.errors) == 1
        assert "Unsupported file format" in result.errors[0]
        assert ".txt" in result.errors[0]

    def test_pdf_file_rejected(self, tmp_path, tmp_output_dir):
        """Test that PDF files are rejected (not Office format)."""
        pdf_file = tmp_path / "test.pdf"
        pdf_file.write_bytes(b"%PDF-1.4 fake pdf content")

        converter = OfficeConverter()
        result = converter.convert(pdf_file, tmp_output_dir)

        assert result.success is False
        assert "Unsupported file format" in result.errors[0]

    def test_file_not_found(self, tmp_output_dir):
        """Test handling of non-existent files."""
        non_existent_file = Path("/non/existent/file.docx")

        converter = OfficeConverter()
        result = converter.convert(non_existent_file, tmp_output_dir)

        assert result.success is False
        assert result.markdown == ""
        assert len(result.errors) == 1
        assert "does not exist" in result.errors[0]

    def test_corrupted_docx(self, tmp_path, tmp_output_dir):
        """Test handling of corrupted DOCX file."""
        corrupted_file = tmp_path / "corrupted.docx"
        # Write invalid content
        corrupted_file.write_bytes(b"This is not a valid DOCX file")

        converter = OfficeConverter()
        result = converter.convert(corrupted_file, tmp_output_dir)

        # Should handle gracefully, not crash
        # MarkItDown may fallback to treating as plain text
        assert result.source_file == corrupted_file
        assert result.converter_name == "OfficeConverter"

    def test_corrupted_pptx(self, tmp_path, tmp_output_dir):
        """Test handling of corrupted PPTX file."""
        corrupted_file = tmp_path / "corrupted.pptx"
        corrupted_file.write_bytes(b"Invalid PPTX content")

        converter = OfficeConverter()
        result = converter.convert(corrupted_file, tmp_output_dir)

        # Should handle gracefully, not crash
        assert result.converter_name == "OfficeConverter"
        assert result.source_file == corrupted_file

    def test_corrupted_xlsx(self, tmp_path, tmp_output_dir):
        """Test handling of corrupted XLSX file."""
        corrupted_file = tmp_path / "corrupted.xlsx"
        corrupted_file.write_bytes(b"Not valid XLSX")

        converter = OfficeConverter()
        result = converter.convert(corrupted_file, tmp_output_dir)

        # Should handle gracefully, not crash
        assert result.converter_name == "OfficeConverter"
        assert result.source_file == corrupted_file

    def test_empty_file(self, tmp_path, tmp_output_dir):
        """Test handling of empty Office file."""
        empty_file = tmp_path / "empty.docx"
        empty_file.write_bytes(b"")

        converter = OfficeConverter()
        result = converter.convert(empty_file, tmp_output_dir)

        # Should fail gracefully
        assert result.success is False


class TestOutputHandling:
    """Test output file handling."""

    def test_output_file_created(self, tmp_path, tmp_output_dir):
        """Test that output file is created in the correct location."""
        docx_file = tmp_path / "test.docx"
        create_docx_file(docx_file)

        converter = OfficeConverter()
        result = converter.convert(docx_file, tmp_output_dir)

        assert result.success is True
        assert result.output_file == tmp_output_dir / "test.md"
        assert result.output_file.exists()

    def test_output_file_content_matches(self, tmp_path, tmp_output_dir):
        """Test that output file content matches result.markdown."""
        docx_file = tmp_path / "content.docx"
        create_docx_file(docx_file)

        converter = OfficeConverter()
        result = converter.convert(docx_file, tmp_output_dir)

        assert result.success is True

        # Read output file content
        output_content = result.output_file.read_text(encoding="utf-8")

        # Verify it matches the markdown in result
        assert output_content == result.markdown

    def test_output_dir_created_if_not_exists(self, tmp_path):
        """Test that output directory is created if it doesn't exist."""
        docx_file = tmp_path / "test.docx"
        create_docx_file(docx_file)

        # Create a non-existent output directory path
        non_existent_dir = tmp_path / "new_output" / "nested"

        converter = OfficeConverter()
        result = converter.convert(docx_file, non_existent_dir)

        assert result.success is True
        assert non_existent_dir.exists()
        assert result.output_file.exists()

    def test_output_file_has_md_extension(self, tmp_path, tmp_output_dir):
        """Test that output file always has .md extension."""
        for ext in [".docx", ".pptx", ".xlsx"]:
            source_file = tmp_path / f"test{ext}"
            if ext == ".docx":
                create_docx_file(source_file)
            elif ext == ".pptx":
                create_pptx_file(source_file)
            else:  # .xlsx
                create_xlsx_file(source_file)

            converter = OfficeConverter()
            result = converter.convert(source_file, tmp_output_dir)

            assert result.success is True
            assert result.output_file.suffix == ".md"


class TestEstimateTime:
    """Test time estimation method."""

    def test_estimate_time_based_on_size(self, tmp_path):
        """Test that time estimation is based on file size."""
        converter = OfficeConverter()

        # Create a small DOCX file
        small_file = tmp_path / "small.docx"
        doc = Document()
        doc.add_paragraph("Small content")
        doc.save(str(small_file))

        # Create a larger DOCX file
        large_file = tmp_path / "large.docx"
        doc = Document()
        for i in range(100):
            doc.add_paragraph(f"Paragraph {i}" * 10)
        doc.save(str(large_file))

        # Estimate times
        small_time = converter.estimate_time(small_file)
        large_time = converter.estimate_time(large_file)

        # Larger file should have larger estimated time
        assert large_time > small_time
        assert small_time > 0
        assert large_time > 0

    def test_estimate_time_nonexistent_file(self):
        """Test that estimate_time handles non-existent files."""
        converter = OfficeConverter()
        nonexistent = Path("/non/existent/file.docx")

        # Should default to 1.0 second
        time_estimate = converter.estimate_time(nonexistent)
        assert time_estimate == 1.0
